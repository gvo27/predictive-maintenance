from pptx import Presentation
from pptx.util import Inches

# Model results
report = {
    "accuracy": 0.985,
    "roc_auc": 0.8991659117210377,
    "classification": {
        "HDF": {"precision": 1.0, "recall": 1.0, "f1": 1.0, "support": 23},
        "NONE": {"precision": 0.985, "recall": 0.999, "f1": 0.992, "support": 1930},
        "OSF": {"precision": 1.0, "recall": 0.5, "f1": 0.667, "support": 16},
        "PWF": {"precision": 0.909, "recall": 0.556, "f1": 0.690, "support": 18},
        "RNF": {"precision": 0.0, "recall": 0.0, "f1": 0.0, "support": 4},
        "TWF": {"precision": 0.0, "recall": 0.0, "f1": 0.0, "support": 9}
    },
    "features": {
        "Torque": 0.242,
        "Tool wear": 0.242,
        "Temp_Delta": 0.171,
        "Rotational speed": 0.168,
        "Air temperature": 0.077,
        "Process temperature": 0.072,
        "Type_L": 0.012,
        "Type_M": 0.009,
        "Type_H": 0.005
    }
}

# Create presentation
prs = Presentation()

# ========== Slide 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Predictive Maintenance — Machine Failure Classifier"
subtitle.text = "UCI AI4I 2020 Dataset"

# ========== Slide 2a: Executive Summary ==========
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Executive Summary"
content = slide.placeholders[1]
content.text = (
    f"Accuracy: {report['accuracy']*100:.1f}%\n"
    f"ROC AUC (macro, OVR): {report['roc_auc']:.3f}\n"
    f"Strong performance in detecting HDF & NONE classes\n"
    f"Weaker detection for rare failures (TWF, RNF)\n\n"
)

# ========== Slide 2b: Executive Summary Cont. ==========
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Executive Summary"
content = slide.placeholders[1]
content.text = (
    "Business Relevance:\n"
    "\t\tSupports Quality Clinic triage with automated failure-type predictions\n"
    "\t\tProvides evidence-based root cause analysis (RCCA)\n"
    "\t\tEnables Supplier Quality & Production to act quickly on discrepancies"
)

# ========== Slide 3: Classification Results ==========
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Classification Performance"

rows = len(report["classification"]) + 1
cols = 5
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table

# headers
headers = ["Class", "Precision", "Recall", "F1-Score", "Support"]
for i, h in enumerate(headers):
    table.cell(0, i).text = h

# data rows
for r, (cls, vals) in enumerate(report["classification"].items(), start=1):
    table.cell(r, 0).text = cls
    table.cell(r, 1).text = f"{vals['precision']:.2f}"
    table.cell(r, 2).text = f"{vals['recall']:.2f}"
    table.cell(r, 3).text = f"{vals['f1']:.2f}"
    table.cell(r, 4).text = str(int(vals['support']))

# ========== Slide 4: Key Drivers ==========
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Top Predictive Features"
content = slide.placeholders[1]
sorted_feats = sorted(report["features"].items(), key=lambda x: x[1], reverse=True)
top_feats_text = "\n".join([f"{i+1}. {name} ({val*100:.1f}%)" for i,(name,val) in enumerate(sorted_feats[:6])])
content.text = (
    "Model identified key drivers of machine failure:\n\n"
    f"{top_feats_text}\n\n"
)

# ========== Slide 5: Business Impact ==========
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Business Significance"
content = slide.placeholders[1]
content.text = (
    "Faster NCR triage — predicted failure type pre-populates disposition options\n"
    "Better traceability — every prediction logged with input evidence\n"
    "Supplier Quality alignment — Torque/Tool Wear insights guide supplier discussions\n"
    "Continuous improvement — feature drivers feed back into preventive maintenance plans\n"
    "Safety first — predictive insights prevent critical machine breakdowns on the floor"
)

# ========== Slide 6: Next Steps ==========
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Next Steps"
content = slide.placeholders[1]
content.text = (
    "Address rare class underperformance (RNF, TWF) using:\n"
    "\t\tSynthetic oversampling (SMOTE)\n"
    "\t\tAdditional labeled data collection\n\n"
    "Deploy classifier into production workflow:\n"
    "\t\tIntegrate with NCR tracking systems\n"
    "\t\tAutomate failure disposition suggestions\n"
)

# Save
prs.save("Quality_Model_Report.pptx")
print("Presentation saved as Quality_Model_Report.pptx")
